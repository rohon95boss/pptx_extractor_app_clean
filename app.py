import io
import json
import os
from typing import List, Tuple

import streamlit as st
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.table import _Cell

# --- .env support (local) ---
from dotenv import load_dotenv
load_dotenv()  # reads OPENAI_API_KEY from .env if present

# --- OpenAI client (works local + Streamlit Cloud/HF Spaces) ---
from openai import OpenAI

# Toggle to see RAW model output for debugging
DEBUG = False  # set True temporarily if you want to inspect raw JSON text


# =========================
# Auth + client management
# =========================
def get_client() -> OpenAI:
    """
    Prefer Streamlit Secrets in the cloud; fall back to local .env in dev.
    Strips hidden whitespace and validates basic format to avoid 401s.
    """
    api_key = (st.secrets.get("OPENAI_API_KEY") if "OPENAI_API_KEY" in st.secrets else None) or os.getenv("OPENAI_API_KEY")
    if not api_key:
        st.error(
            "Missing **OPENAI_API_KEY**.\n\n"
            "• Local dev: create a `.env` with `OPENAI_API_KEY=...`\n"
            "• Streamlit Cloud / HF Spaces: Settings → Secrets/Variables → add `OPENAI_API_KEY`"
        )
        st.stop()

    api_key = api_key.strip()  # remove trailing spaces/newlines
    if not api_key.startswith(("sk-", "sk-proj-")):
        st.error("Your OPENAI_API_KEY looks malformed (check quotes/spaces).")
        st.stop()

    return OpenAI(api_key=api_key)


def _auth_diag():
    """Sidebar diagnostics to pinpoint auth issues quickly."""
    st.sidebar.markdown("### 🔐 Auth Diagnostics")
    src = "st.secrets" if "OPENAI_API_KEY" in st.secrets else ("os.getenv" if os.getenv("OPENAI_API_KEY") else "none")
    st.sidebar.write(f"Source: **{src}**")

    # safe preview (no leak)
    val = (st.secrets.get("OPENAI_API_KEY") if "OPENAI_API_KEY" in st.secrets else os.getenv("OPENAI_API_KEY") or "").strip()
    st.sidebar.write(f"Key present: **{'yes' if val else 'no'}**")
    if val:
        st.sidebar.write(f"Key preview: `{val[:6]}…{val[-4:]}`")

    # live probe
    try:
        client = OpenAI(api_key=val) if val else None
        if client:
            client.models.list()  # cheap call
            st.sidebar.success("Auth probe: OK (models.list)")
        else:
            st.sidebar.error("Auth probe failed: no key available")
    except Exception as e:
        st.sidebar.error(f"Auth probe failed: {e}")

    # minimal smoke test button
    if st.sidebar.button("🧪 JSON Smoke Test"):
        try:
            client = get_client()
            r = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[
                    {"role": "system", "content": "Return ONLY JSON: {\"ok\": true}"},
                    {"role": "user", "content": "Say hello as JSON."},
                ],
                temperature=0,
                response_format={"type": "json_object"},
            )
            st.sidebar.success(f"OK: {r.choices[0].message.content}")
        except Exception as e:
            st.sidebar.error(f"Smoke failed: {e}")


# =========================
# Business logic + schemas
# =========================
ALLOWED_FUNCTIONS = [
    "COO", "Technology", "Compliance", "Risk", "Operations", "Finance", "Treasury",
    "Front Office", "Middle Office", "Back Office", "Data", "Regulatory", "Legal", "Internal Audit"
]

class AnalysisOutput(BaseModel):
    case_name: str = Field('', description="Name of the case/study/project")
    category: str = Field('', description="High-level category")
    function: str = Field('', description="Office/organizational unit (e.g., COO, Technology, Compliance)")
    hashtags: List[str] = Field(default_factory=list, description="Exactly 3 social hashtags (no #)")
    challenge: str = Field('', description="The challenge/pain/problem")
    solution: str = Field('', description="The solution/approach")
    results: str = Field('', description="The impact/outcomes/metrics")
    business_processes: List[str] = Field(default_factory=list, description="Exactly 5 process bullets")


SYSTEM_PROMPT = """You are a precise business-case information extractor for management consulting artifacts.
Return ONLY valid JSON matching this schema:

{
  "case_name": "string",
  "category": "string",
  "function": "string",
  "hashtags": ["string","string","string"],
  "challenge": "string",
  "solution": "string",
  "results": "string",
  "business_processes": ["string","string","string","string","string"]
}

Strict rules:
1) Anonymize the client. Never output a real client or bank name. Use descriptors like:
   "Top US bank", "Global investment bank", "Tier-1 broker-dealer", "Fortune 100 insurer",
   "Leading payments provider", "Major asset manager", "Top card issuer", etc.
2) Always mention our firm explicitly as "BIP" in the narrative (challenge/solution/results).
3) Case name: concise, professional, and anonymized. Preferred pattern:
   "BIP Case Study – <Anonymized Client Type> – <Project Theme>"
4) Category: a business-friendly label (e.g., "Regulatory Reporting", "Risk", "Operations", "Technology",
   "Data", "Compliance", "Change", "Trading").
5) Function: MUST be the office / organizational unit (choose best fit from the allowed list).
6) Hashtags: exactly 3, no leading '#', short social-style, business-relevant.
7) Challenge, Solution, Results: ~80–140 words each, integrate salient info from source, no PII.
8) Business processes: exactly 5 short noun phrases in Title Case from typical consulting activities.
9) Do NOT invent facts outside the source; stay generic if details are missing.
10) Output MUST be ONLY the JSON object. No extra text, no markdown.
"""


# =========================
# PPTX Extraction
# =========================
def extract_text_from_pptx_with_progress(file_bytes: bytes) -> Tuple[str, dict]:
    prs = Presentation(io.BytesIO(file_bytes))
    num_slides = len(prs.slides)

    chunks: List[str] = []
    table_rows = 0
    notes_lines = 0
    text_shapes = 0

    st.write(f"Found **{num_slides}** slides.")
    pbar = st.progress(0)
    prog_txt = st.empty()

    def shape_text(shape) -> List[str]:
        nonlocal table_rows, text_shapes
        out = []

        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            text_shapes += 1
            for p in shape.text_frame.paragraphs:
                out.append("".join(run.text for run in p.runs).strip())

        if hasattr(shape, "has_table") and shape.has_table:
            table = shape.table
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    if isinstance(cell, _Cell):
                        cell_txt = (cell.text or "").strip()
                        if cell_txt:
                            cells.append(cell_txt)
                if cells:
                    out.append(" | ".join(cells))
                    table_rows += 1

        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            for s in shape.shapes:
                out.extend(shape_text(s))

        return [x for x in out if x]

    for idx, slide in enumerate(prs.slides, start=1):
        try:
            if slide.shapes.title and slide.shapes.title.text:
                chunks.append(slide.shapes.title.text.strip())
        except Exception:
            pass

        for s in slide.shapes:
            chunks.extend(shape_text(s))

        if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                for p in notes_frame.paragraphs:
                    note_line = "".join(run.text for run in p.runs).strip()
                    if note_line:
                        chunks.append(f"[Notes] {note_line}")
                        notes_lines += 1

        pbar.progress(int(idx / max(1, num_slides) * 100))
        prog_txt.write(f"Extracting slide {idx}/{num_slides}…")

    cleaned = [c for c in (t.strip() for t in chunks) if c]
    final_text = "\n".join(cleaned)

    stats = {
        "slides": num_slides,
        "text_shapes": text_shapes,
        "table_rows": table_rows,
        "notes_lines": notes_lines,
        "chunks": len(cleaned),
        "characters": len(final_text),
    }
    return final_text, stats


# =========================
# Post-processing helpers
# =========================
def _map_function_to_office(func_value: str, raw_text: str) -> str:
    f = (func_value or "").strip()
    for canon in ALLOWED_FUNCTIONS:
        if f and f.lower() == canon.lower():
            return canon

    text_l = (raw_text or "").lower() + " " + (f.lower() if f else "")
    mapping = [
        ("chief operating officer", "COO"),
        ("coo", "COO"),
        ("operations", "Operations"),
        ("ops", "Operations"),
        ("technology", "Technology"),
        ("tech", "Technology"),
        ("it ", "Technology"),
        ("compliance", "Compliance"),
        ("risk", "Risk"),
        ("finance", "Finance"),
        ("treasury", "Treasury"),
        ("front office", "Front Office"),
        ("trading", "Front Office"),
        ("sales", "Front Office"),
        ("middle office", "Middle Office"),
        ("controls", "Middle Office"),
        ("back office", "Back Office"),
        ("settlement", "Back Office"),
        ("data", "Data"),
        ("regulatory", "Regulatory"),
        ("legal", "Legal"),
        ("internal audit", "Internal Audit"),
        ("audit", "Internal Audit"),
    ]
    for kw, label in mapping:
        if kw in text_l:
            return label
    return "Operations"


def _normalize_output(o: AnalysisOutput, raw_text: str) -> AnalysisOutput:
    if "BIP" not in (o.solution or ""):
        o.solution = (o.solution + " ").strip() + "BIP led the delivery and execution."
    if "BIP" not in (o.results or ""):
        o.results = (o.results + " ").strip() + "BIP enabled measurable business outcomes."

    seen = set()
    cleaned = []
    for h in o.hashtags:
        if not h:
            continue
        h2 = h.strip().lstrip("#").lower().replace(" ", "")[:24]
        if h2 and h2 not in seen:
            cleaned.append(h2); seen.add(h2)
    while len(cleaned) < 3:
        cleaned.append("")
    o.hashtags = cleaned[:3]

    bp_seen = set()
    bps = []
    for bp in o.business_processes:
        if not bp:
            continue
        bp2 = " ".join(bp.split()).title()
        if bp2 not in bp_seen:
            bps.append(bp2); bp_seen.add(bp2)
    while len(bps) < 5:
        bps.append("")
    o.business_processes = bps[:5]

    o.function = _map_function_to_office(o.function, raw_text)

    if not o.case_name.strip():
        o.case_name = "BIP Case Study – Top US Bank – Consulting Engagement"

    return o


# =========================
# OpenAI call (JSON mode)
# =========================
def analyze_text_with_openai(raw_text: str) -> AnalysisOutput:
    client = get_client()

    user_prompt = (
        "Content from a PowerPoint presentation follows between <<< and >>>. "
        "Extract and summarize into the exact JSON schema described by the system. "
        "Apply anonymization and BIP rules as specified. "
        "For Challenge, Solution, Results: produce elaborated, business-context paragraphs (~80–140 words each) "
        "that integrate salient, defensible details from the provided content without exposing real client names. "
        "ONLY output JSON. No markdown or extra prose.\n\n"
        f"<<<\n{raw_text}\n>>>"
    )

    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini",  # supports JSON mode; reliable and cheaper
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0,
            max_tokens=900,
            response_format={"type": "json_object"},  # force strict JSON
        )
    except Exception as e:
        msg = str(e)
        if "401" in msg or "invalid_api_key" in msg.lower():
            raise RuntimeError(
                "OpenAI 401 invalid_api_key. Check Secrets/ENV (exact key, no spaces), "
                "SDK version, and remove any project-key origin/IP restrictions."
            )
        raise RuntimeError(f"OpenAI request failed: {e}")

    content = (resp.choices[0].message.content or "").strip()

    if DEBUG:
        with st.expander("RAW model output (debug)"):
            st.code(content, language="json")

    try:
        data = json.loads(content)
    except Exception:
        # Should be rare with JSON mode; show raw for diagnosis if needed
        clipped = content[content.find("{"): content.rfind("}") + 1]
        try:
            data = json.loads(clipped)
        except Exception:
            raise ValueError(f"Model did not return valid JSON.\n\nRAW:\n{content}")

    out = AnalysisOutput(**data)

    if len(out.hashtags) != 3:
        out.hashtags = (out.hashtags + ["", "", ""])[:3]
    if len(out.business_processes) != 5:
        out.business_processes = (out.business_processes + [""] * 5)[:5]

    out = _normalize_output(out, raw_text)
    return out


# =========================
# UI
# =========================
st.set_page_config(page_title="PPT → Structured Case Analyzer", page_icon="📊", layout="centered")
_auth_diag()  # remove once green

st.title("📊 PowerPoint Text Extraction → Case Analyzer")

with st.expander("How it works", expanded=False):
    st.markdown(
        "1) Upload a `.pptx`\n"
        "2) We extract all text (slides, tables, groups, notes)\n"
        "3) You confirm or paste any missing text\n"
        "4) Click **Analyze** → get structured fields\n"
    )

uploaded = st.file_uploader("Upload a PowerPoint (.pptx)", type=["pptx"])

if uploaded:
    st.success("File received.")
    ppt_bytes = uploaded.read()

    with st.status("Extracting text from slides…", expanded=True) as s:
        extracted, stats = extract_text_from_pptx_with_progress(ppt_bytes)
        s.update(label="Extraction complete ✅", state="complete")

    m1, m2, m3, m4 = st.columns(4)
    m1.metric("Slides", stats["slides"])
    m2.metric("Text Shapes", stats["text_shapes"])
    m3.metric("Table Rows", stats["table_rows"])
    m4.metric("Notes Lines", stats["notes_lines"])
    st.caption(f"Detected **{stats['chunks']}** text chunks · **{stats['characters']}** characters")

    st.subheader("Step 1 — Review & confirm the extracted text")
    st.caption("If anything is missing (e.g., text baked into images), add it below before analysis.")
    user_text = st.text_area(
        "All extracted text (editable):",
        extracted,
        height=350,
        placeholder="All slide text will appear here…",
    )

    c1, c2 = st.columns(2)
    with c1:
        confirm = st.button("✅ Yes, looks complete")
    with c2:
        add_more = st.button("✍️ I added missing text")

    proceed = confirm or add_more

    if proceed:
        st.subheader("Step 2 — Run AI analysis")
        with st.status("Analyzing with OpenAI…", expanded=True) as s:
            st.write("Packing structured prompt…")
            try:
                result = analyze_text_with_openai(user_text)
                st.write("Parsing JSON…")
                st.write("Validating and normalizing output…")
            except Exception as e:
                s.update(label="Analysis failed ❌", state="error")
                st.error(f"OpenAI call or JSON parsing failed: {e}")
            else:
                s.update(label="Analysis complete ✅", state="complete")
                st.success("Done!")

                st.markdown("### JSON Result")
                st.json(json.loads(result.model_dump_json()))

                st.markdown("### Summary")
                st.write(f"**Case Name:** {result.case_name or '-'}")
                st.write(f"**Category:** {result.category or '-'}")
                st.write(f"**Function:** {result.function or '-'}")
                st.write("**Hashtags:** " + (", ".join([h for h in result.hashtags if h]) or "-"))
                st.write("**Challenge:**")
                st.write(result.challenge or "-")
                st.write("**Solution:**")
                st.write(result.solution or "-")
                st.write("**Results:**")
                st.write(result.results or "-")
                st.write("**Business Processes (5):**")
                for i, bp in enumerate(result.business_processes, start=1):
                    st.write(f"{i}. {bp or '-'}")

                st.download_button(
                    "⬇️ Download JSON",
                    data=result.model_dump_json(indent=2),
                    file_name="case_analysis.json",
                    mime="application/json",
                )
else:
    st.info("Upload a `.pptx` to begin.")
